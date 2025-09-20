import os
import json
import uuid
import re
from pathlib import Path
from typing import List, Dict, Any, Tuple, Optional
from io import StringIO

from pdfminer.high_level import extract_text_to_fp
from pdfminer.layout import LAParams

# Simple text cleaning function to replace unstructured's clean_extra_whitespace
def clean_extra_whitespace(text):
    """Clean extra whitespace from text."""
    if not text:
        return ""
    # Replace multiple whitespace with single space
    text = re.sub(r'\s+', ' ', text)
    # Remove whitespace at the beginning and end
    return text.strip()

from app.config.settings import settings


class DocumentProcessor:
    """
    Utility class for processing documents (PDF, DOCX, PPTX, etc.).
    Handles parsing, text extraction, and chunking.
    """
    
    def __init__(self, file_path: str, file_type: str, use_ultra_fast_processing: bool = False):
        """
        Initialize the document processor.
        
        Args:
            file_path: Path to the document file.
            file_type: Type of the document (pdf, docx, pptx, etc.).
            use_ultra_fast_processing: Whether to use ultra-fast processing for small files.
        """
        self.file_path = file_path
        self.file_type = file_type.lower()
        self.text_content = ""
        self.page_count = 0
        self.pages = []
        self.use_ultra_fast_processing = use_ultra_fast_processing
    
    def process(self) -> bool:
        """
        Process the document based on its type.
        
        Returns:
            bool: True if processing was successful, False otherwise.
        """
        try:
            print(f"🔍 DocumentProcessor: Processing {self.file_type} file at {self.file_path}")
            
            if self.file_type == "pdf":
                print(f"📄 Processing PDF...")
                return self.process_pdf()
            elif self.file_type == "docx":
                print(f"📄 Processing DOCX...")
                return self._process_docx()
            elif self.file_type == "pptx":
                print(f"📄 Processing PPTX...")
                return self._process_pptx()
            elif self.file_type in ["txt", "md"]:
                print(f"📄 Processing text file...")
                return self._process_text()
            else:
                print(f"❌ Unsupported file type: {self.file_type}")
                return False
        except Exception as e:
            import traceback
            print(f"❌ Error processing document: {e}")
            print(f"❌ DocumentProcessor error details: {traceback.format_exc()}")
            return False
    
    def process_pdf(self) -> bool:
        """
        Process a PDF file and extract text from each page.
        
        Returns:
            bool: True if processing was successful, False otherwise.
        """
        try:
            import os
            
            # Ultra-fast path for small files when flag is set
            if self.use_ultra_fast_processing:
                return self._process_pdf_ultra_fast()
            
            from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
            from pdfminer.converter import TextConverter
            from pdfminer.layout import LAParams
            from pdfminer.pdfpage import PDFPage
            from pdfminer.pdfdocument import PDFSyntaxError
            from io import StringIO
            
            # First, validate the PDF and count pages
            try:
                with open(self.file_path, 'rb') as file:
                    pages = list(PDFPage.get_pages(file))
                    if not pages:
                        print("No pages found in PDF")
                        return False
                    
                    # Rewind the file for later processing
                    file.seek(0)
            except PDFSyntaxError as e:
                print(f"Invalid PDF format: {e}")
                return False
            
            # Extract text from each page
            from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
            from pdfminer.converter import TextConverter
            from pdfminer.pdfpage import PDFPage
            
            import threading
            import time
            from concurrent.futures import ThreadPoolExecutor, TimeoutError
            
            # Reset page list
            self.pages = []
            max_processing_time = 30  # Maximum processing time in seconds
            
            def process_pdf_with_timeout():
                with open(self.file_path, 'rb') as file:
                    # Process each page individually with a timeout
                    for page_num, page in enumerate(PDFPage.get_pages(file)):
                        output_string = StringIO()
                        resource_manager = PDFResourceManager()
                        device = TextConverter(resource_manager, output_string, laparams=LAParams())
                        interpreter = PDFPageInterpreter(resource_manager, device)
                        
                        # Process this page
                        interpreter.process_page(page)
                        device.close()
                        
                        # Get and clean the text
                        text = output_string.getvalue()
                        text = clean_extra_whitespace(text)
                        
                        # Add to pages collection
                        self.pages.append({
                            "page_number": page_num + 1,
                            "content": text
                        })
                return True
                
            # Use a thread pool to enforce timeout
            with ThreadPoolExecutor(max_workers=1) as executor:
                future = executor.submit(process_pdf_with_timeout)
                try:
                    result = future.result(timeout=max_processing_time)
                    if not result:
                        print(f"PDF processing failed or returned unexpected result")
                        return False
                except TimeoutError:
                    print(f"PDF processing timed out after {max_processing_time} seconds")
                    return False
                except Exception as e:
                    print(f"Error during PDF processing: {e}")
                    return False
            
            # Combine all text
            self.text_content = "\n\n".join([page["content"] for page in self.pages])
            return True
        except ImportError as e:
            print(f"Missing PDF processing dependency: {e}")
            print("Please install required dependencies: pip install pdfminer.six")
            return False
        except FileNotFoundError:
            print(f"PDF file not found or not accessible: {self.file_path}")
            return False
        except Exception as e:
            import traceback
            print(f"Error processing PDF: {e}")
            print(f"PDF processing stack trace: {traceback.format_exc()}")
            return False
    
    def _process_docx(self) -> bool:
        """
        Process a DOCX document.
        
        Returns:
            bool: True if processing was successful, False otherwise.
        """
        try:
            print(f"📄 Starting DOCX processing...")
            from docx import Document as DocxDocument
            
            print(f"📖 Loading DOCX document...")
            doc = DocxDocument(self.file_path)
            full_text = []
            
            # Extract text from paragraphs
            print(f"📝 Extracting paragraphs...")
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    full_text.append(paragraph.text.strip())
            
            # Extract text from tables
            print(f"📊 Extracting tables...")
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        full_text.append(" | ".join(row_text))
            
            self.text_content = "\n".join(full_text)
            self.page_count = 1  # DOCX doesn't have explicit pages, so we use 1
            self.pages = [{"page_number": 1, "content": self.text_content}]
            
            print(f"✅ DOCX processing complete. Extracted {len(self.text_content)} characters")
            return True
        except ImportError as e:
            print(f"❌ Missing DOCX processing dependency: {e}")
            print("Please install required dependency: pip install python-docx")
            return False
        except FileNotFoundError:
            print(f"❌ DOCX file not found or not accessible: {self.file_path}")
            return False
        except Exception as e:
            import traceback
            print(f"❌ Error processing DOCX: {e}")
            print(f"❌ DOCX processing stack trace: {traceback.format_exc()}")
            return False
    
    def _process_pptx(self) -> bool:
        """
        Process a PPTX document.
        
        Returns:
            bool: True if processing was successful, False otherwise.
        """
        try:
            # Simple implementation for PPTX processing
            import pptx
            
            # Try to import pptx, if not available, provide helpful error
            try:
                presentation = pptx.Presentation(self.file_path)
            except ImportError:
                print("python-pptx not installed. Install with: pip install python-pptx")
                return False
            
            # Process each slide
            slide_num = 1
            for slide in presentation.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(clean_extra_whitespace(shape.text))
                
                if slide_text:  # Only add slides with text content
                    self.pages.append({
                        "page_number": slide_num,
                        "content": "\n".join(slide_text)
                    })
                    slide_num += 1
            
            # Combine all text
            self.text_content = "\n\n".join([page["content"] for page in self.pages])
            self.page_count = len(self.pages)
            
            return True
        except ImportError as e:
            print(f"Missing PPTX processing dependency: {e}")
            print("Please install required dependency: pip install python-pptx")
            return False
        except FileNotFoundError:
            print(f"PPTX file not found or not accessible: {self.file_path}")
            return False
        except Exception as e:
            import traceback
            print(f"Error processing PPTX: {e}")
            print(f"PPTX processing stack trace: {traceback.format_exc()}")
            return False
    
    def chunk_document(self, chunk_size: int = None, chunk_overlap: int = None) -> List[Dict[str, Any]]:
        """
        Split the document into chunks for embedding.
        
        Args:
            chunk_size: Size of each chunk in characters.
            chunk_overlap: Overlap between chunks in characters.
            
        Returns:
            List[Dict[str, Any]]: List of document chunks with metadata.
        """
        print(f"✂️ Starting chunking process...")
        if chunk_size is None:
            chunk_size = settings.chunk_size
        if chunk_overlap is None:
            chunk_overlap = settings.chunk_overlap
        
        print(f"✂️ Chunk size: {chunk_size}, Overlap: {chunk_overlap}")
        chunks = []
        
        # Process each page separately to maintain page references
        print(f"✂️ Processing {len(self.pages)} pages...")
        for i, page in enumerate(self.pages):
            # Use consistent key name for page number
            page_num = page.get("page_number", page.get("page_num", i + 1))
            text = page["content"]
            
            # Skip empty pages
            if not text.strip():
                continue
            
            # If the page is smaller than chunk_size, use it as a single chunk
            if len(text) <= chunk_size:
                chunks.append({
                    "content": text,
                    "metadata": {
                        "page_number": page_num,
                        "chunk_index": len(chunks),
                        "source": os.path.basename(self.file_path)
                    }
                })
                continue
            
            # Split the page into chunks efficiently
            start = 0
            max_iterations = len(text) // max(1, chunk_size - chunk_overlap) + 10  # Safety limit
            iteration_count = 0
            
            while start < len(text) and iteration_count < max_iterations:
                iteration_count += 1
                # Find the end of the chunk
                end = start + chunk_size
                
                # If we're not at the end of the text, try to find a good breaking point
                if end < len(text):
                    # Try to find a period, question mark, or exclamation mark followed by a space or newline
                    for i in range(end - 1, start + chunk_size // 2, -1):
                        if i < len(text) and text[i] in ['.', '!', '?', '\n'] and (i + 1 == len(text) or text[i + 1].isspace()):
                            end = i + 1
                            break
                else:
                    end = len(text)
                
                # Create the chunk
                chunk_text = text[start:end].strip()
                if chunk_text:  # Only add non-empty chunks
                    chunks.append({
                        "content": chunk_text,
                        "metadata": {
                            "page_number": page_num,
                            "chunk_index": len(chunks),
                            "source": os.path.basename(self.file_path)
                        }
                    })
                
                # Move to the next chunk with overlap, ensuring we always advance
                next_start = end - chunk_overlap
                if next_start <= start:  # Safety check to prevent infinite loops
                    next_start = start + max(1, chunk_size // 2)  # Force advancement
                start = next_start
        
        print(f"✂️ Chunking complete! Created {len(chunks)} chunks")
        return chunks
    
    def _process_pdf_ultra_fast(self) -> bool:
        """
        Ultra-fast PDF processing for very small files (under 10KB).
        Uses minimal processing to extract text quickly.
        """
        try:
            import PyPDF2
            from io import StringIO
            
            # Try PyPDF2 first for speed
            try:
                with open(self.file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    
                    # Reset pages list
                    self.pages = []
                    
                    # Extract text from each page quickly
                    for page_num, page in enumerate(pdf_reader.pages):
                        try:
                            text = page.extract_text()
                            if text.strip():  # Only add non-empty pages
                                self.pages.append({
                                    "page_number": page_num + 1,
                                    "content": text.strip()
                                })
                        except Exception as e:
                            print(f"Error extracting text from page {page_num + 1}: {e}")
                            continue
                    
                    # Combine all text
                    self.text_content = "\n\n".join([page["content"] for page in self.pages])
                    print(f"Ultra-fast processing complete: {len(self.pages)} pages, {len(self.text_content)} characters")
                    return True
                    
            except Exception as e:
                print(f"PyPDF2 failed, falling back to pdfplumber: {e}")
                
            # Fallback to pdfplumber for better text extraction
            try:
                import pdfplumber
                
                # Reset pages list
                self.pages = []
                
                with pdfplumber.open(self.file_path) as pdf:
                    for page_num, page in enumerate(pdf.pages):
                        try:
                            text = page.extract_text()
                            if text and text.strip():
                                self.pages.append({
                                    "page_number": page_num + 1,
                                    "content": text.strip()
                                })
                        except Exception as e:
                            print(f"Error extracting text from page {page_num + 1}: {e}")
                            continue
                
                # Combine all text
                self.text_content = "\n\n".join([page["content"] for page in self.pages])
                print(f"Ultra-fast processing with pdfplumber complete: {len(self.pages)} pages")
                return True
                
            except ImportError:
                print("pdfplumber not available, falling back to standard processing")
                return False
            except Exception as e:
                print(f"pdfplumber processing failed: {e}")
                return False
                
        except Exception as e:
            print(f"Ultra-fast processing failed: {e}")
            return False
    
    def _process_text(self) -> bool:
        """
        Process a plain text file.
        
        Returns:
            bool: True if processing was successful, False otherwise.
        """
        try:
            print(f"📄 Starting text file processing...")
            
            with open(self.file_path, 'r', encoding='utf-8') as file:
                text_content = file.read()
            
            self.text_content = text_content
            self.page_count = 1  # Text files are treated as single page
            self.pages = [{"page_number": 1, "content": text_content}]
            
            print(f"✅ Text processing complete. Extracted {len(self.text_content)} characters")
            return True
        except Exception as e:
            import traceback
            print(f"❌ Error processing text file: {e}")
            print(f"❌ Text processing stack trace: {traceback.format_exc()}")
            return False
    
    def get_page_count(self) -> int:
        """
        Get the number of pages in the document.
        
        Returns:
            int: Number of pages.
        """
        return self.page_count
    
    def get_text_content(self) -> str:
        """
        Get the full text content of the document.
        
        Returns:
            str: Full text content.
        """
        return self.text_content
    
    def get_pages(self) -> List[Dict[str, Any]]:
        """
        Get the pages of the document.
        
        Returns:
            List[Dict[str, Any]]: List of pages with content.
        """
        return self.pages
